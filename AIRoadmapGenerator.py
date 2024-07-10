import streamlit as st
import pandas as pd
import datetime
from pptx import Presentation
from pptx.util import Inches

# Define the use cases
use_cases = [
    {"Category": "Product", "Use Case": "Personalized product recommendations", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Product", "Use Case": "Dynamic pricing optimization", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Product", "Use Case": "Product lifecycle management", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Product", "Use Case": "Customer churn prediction", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Marketing", "Use Case": "Targeted marketing campaigns", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Marketing", "Use Case": "Sentiment analysis on social media", "Audience": "Client Focus", "Impact": "Productivity Improvements"},
    {"Category": "Marketing", "Use Case": "Marketing spend optimization", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "Marketing", "Use Case": "Lead scoring and prioritization", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Sales", "Use Case": "Sales forecasting", "Audience": "Employee Focus", "Impact": "Revenue Increase"},
    {"Category": "Sales", "Use Case": "Customer segmentation", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Sales", "Use Case": "Upsell and cross-sell recommendations", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Sales", "Use Case": "Automated sales reporting", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Customer Service", "Use Case": "Chatbots for customer support", "Audience": "Client Focus", "Impact": "Cost Savings"},
    {"Category": "Customer Service", "Use Case": "Automated email responses", "Audience": "Client Focus", "Impact": "Cost Savings"},
    {"Category": "Customer Service", "Use Case": "Predictive customer service", "Audience": "Client Focus", "Impact": "Productivity Improvements"},
    {"Category": "Customer Service", "Use Case": "Customer feedback analysis", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Controls", "Use Case": "Regulatory compliance automation", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "Controls", "Use Case": "Automated audit processes", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Controls", "Use Case": "Risk management and monitoring", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "Controls", "Use Case": "Policy adherence monitoring", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Risk", "Use Case": "Credit risk assessment", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Risk", "Use Case": "Fraud risk detection", "Audience": "Client Focus", "Impact": "Cost Savings"},
    {"Category": "Risk", "Use Case": "Market risk analysis", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Risk", "Use Case": "Operational risk management", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Fraud", "Use Case": "Real-time fraud detection", "Audience": "Client Focus", "Impact": "Cost Savings"},
    {"Category": "Fraud", "Use Case": "Transaction anomaly detection", "Audience": "Client Focus", "Impact": "Cost Savings"},
    {"Category": "Fraud", "Use Case": "Fraudulent document detection", "Audience": "Client Focus", "Impact": "Cost Savings"},
    {"Category": "Fraud", "Use Case": "Fraud risk scoring", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Operations", "Use Case": "Process automation", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "Operations", "Use Case": "Inventory management", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Operations", "Use Case": "Supply chain optimization", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "Operations", "Use Case": "Workforce scheduling", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "HR", "Use Case": "Employee attrition prediction", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "HR", "Use Case": "Recruitment process automation", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "HR", "Use Case": "Performance management", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "HR", "Use Case": "Learning and development recommendations", "Audience": "Employee Focus", "Impact": "Productivity Improvements"},
    {"Category": "Finance", "Use Case": "Financial forecasting", "Audience": "Employee Focus", "Impact": "Revenue Increase"},
    {"Category": "Finance", "Use Case": "Expense management automation", "Audience": "Employee Focus", "Impact": "Cost Savings"},
    {"Category": "Finance", "Use Case": "Investment portfolio optimization", "Audience": "Client Focus", "Impact": "Revenue Increase"},
    {"Category": "Finance", "Use Case": "Automated financial reporting", "Audience": "Employee Focus", "Impact": "Productivity Improvements"}
]

# Convert the use cases to a DataFrame
use_cases_df = pd.DataFrame(use_cases)

# Initialize session state
if "assigned_use_cases" not in st.session_state:
    st.session_state.assigned_use_cases = []
if "total_budget" not in st.session_state:
    st.session_state.total_budget = 100000  # Default total budget

# Define functions for PPT generation
def generate_ppt(assigned_use_cases):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "AI Roadmap"

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(0.5)

    for use_case in assigned_use_cases:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"{use_case['Use Case']} - {use_case['Category']} - {use_case['Audience']} - {use_case['Impact']} - {use_case['Importance']} - {use_case['Urgency']} - {use_case['Timeline']} - {use_case['Budget']}"
        top += Inches(0.5)
    
    return prs

# UI for Use Case Assignment
st.title("AI Roadmap Generator")

st.subheader("Assign Importance, Urgency, Timelines, and Budgets to Use Cases")
################################

# Filters
department_filter = st.selectbox("Filter by Department", ["All"] + sorted(use_cases_df["Category"].unique()))
audience_filter = st.selectbox("Filter by Audience", ["All", "Client Focus", "Employee Focus"])
impact_filter = st.selectbox("Filter by Impact", ["All", "Revenue Increase", "Cost Savings", "Productivity Improvements"])

filtered_use_cases_df = use_cases_df
if department_filter != "All":
    filtered_use_cases_df = filtered_use_cases_df[filtered_use_cases_df["Category"] == department_filter]
if audience_filter != "All":
    filtered_use_cases_df = filtered_use_cases_df[filtered_use_cases_df["Audience"] == audience_filter]
if impact_filter != "All":
    filtered_use_cases_df = filtered_use_cases_df[filtered_use_cases_df["Impact"] == impact_filter]

# Display filtered use cases
st.write("Filtered Use Cases:")
st.write(filtered_use_cases_df)

# Display all use cases
st.write("All Use Cases:")
st.write(use_cases_df)
################################
# Function to create Gantt chart
def create_gantt_chart(assigned_use_cases):
    if assigned_use_cases:
        df = pd.DataFrame(assigned_use_cases)
        df['Start'] = df['Timeline']
        df['Finish'] = df['Start'] + pd.DateOffset(months=2)  # Fixed duration of 2 months
        fig = px.timeline(df, x_start="Start", x_end="Finish", y="Use Case", color="Category",
                          title="AI Roadmap Gantt Chart", labels={"Use Case": "AI Use Cases"})
        fig.update_yaxes(categoryorder="total ascending")
        st.plotly_chart(fig)

# Assign fixed duration and budget to use cases
selected_use_case = st.selectbox("Select a Use Case", filtered_use_cases_df["Use Case"])
importance = st.slider("Importance", 1, 10)
urgency = st.slider("Urgency", 1, 10)
timeline = datetime.date.today() + datetime.timedelta(days=60)  # Fixed duration of 2 months
budget = 10000  # Fixed budget of $10k

if st.button("Assign Use Case"):
    assigned_budget = sum(use_case["Budget"] for use_case in st.session_state.assigned_use_cases)

    if assigned_budget + budget > st.session_state.total_budget:
        st.error("Assigning this use case would exceed the total budget. Please adjust other use cases or increase the total budget.")
    else:
        assigned_use_case = {
            "Use Case": selected_use_case,
            "Category": filtered_use_cases_df[filtered_use_cases_df["Use Case"] == selected_use_case]["Category"].values[0],
            "Audience": filtered_use_cases_df[filtered_use_cases_df["Use Case"] == selected_use_case]["Audience"].values[0],
            "Impact": filtered_use_cases_df[filtered_use_cases_df["Use Case"] == selected_use_case]["Impact"].values[0],
            "Importance": importance,
            "Urgency": urgency,
            "Timeline": timeline,
            "Budget": budget
        }
        st.session_state.assigned_use_cases.append(assigned_use_case)

# Display assigned use cases
st.write("Assigned Use Cases:")
assigned_use_cases_df = pd.DataFrame(st.session_state.assigned_use_cases)
st.write(assigned_use_cases_df)

# UI for Gantt Chart Visualization (Mockup)
st.subheader("Gantt Chart Visualization (Mockup)")
# Note: In a real implementation, you would use a library to generate a Gantt chart
st.write("Gantt Chart visualization would appear here.")

# Generate PowerPoint
if st.button("Generate PowerPoint"):
    prs = generate_ppt(st.session_state.assigned_use_cases)
    prs.save('/tmp/AI_Roadmap.pptx')
    with open('/tmp/AI_Roadmap.pptx', "rb") as ppt_file:
        st.download_button(
            label="Download PowerPoint",
            data=ppt_file,
            file_name="AI_Roadmap.pptx"
        )

st.sidebar.subheader("Budget Tracker")
total_budget = st.sidebar.number_input("Total Budget", min_value=0)
assigned_budget = sum(use_case["Budget"] for use_case in st.session_state.assigned_use_cases)
remaining_budget = total_budget - assigned_budget

st.sidebar.write(f"Assigned Budget: {assigned_budget}")
st.sidebar.write(f"Remaining Budget: {remaining_budget}")

if remaining_budget < 0:
    st.sidebar.error("Budget exceeded! Adjust your use cases or increase the total budget.")
else:
    st.sidebar.success("Budget is within the limit.")
